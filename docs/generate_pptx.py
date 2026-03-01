#!/usr/bin/env python3
"""Generate NitNab Design Review .pptx presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import datetime

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BG_DARK = RGBColor(0x1D, 0x1D, 0x1F)
BG_CARD = RGBColor(0x2C, 0x2C, 0x2E)
ACCENT = RGBColor(0x00, 0x71, 0xE3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0x98, 0x98, 0x9D)
GREEN = RGBColor(0x30, 0xD1, 0x58)
RED = RGBColor(0xFF, 0x45, 0x3A)
ORANGE = RGBColor(0xFF, 0x9F, 0x0A)
YELLOW = RGBColor(0xFF, 0xD6, 0x0A)

def set_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color=BG_CARD, radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if radius is not None:
        shape.adjustments[0] = radius
    return shape

def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tf

def add_bullet_list(slide, left, top, width, height, items, size=16, color=WHITE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(6)
        p.level = 0
    return tf

# =====================================================================
# SLIDE 1: TITLE
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_bg(slide)

add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
         'NitNab', size=60, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
         'Design Review & Feature Implementation Plan', size=28, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.5),
         'macOS Speech Transcription Application', size=20, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.8),
         f'Design & Engineering Review Team  |  {datetime.date.today().strftime("%B %d, %Y")}  |  Confidential',
         size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 2: AGENDA
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         'Agenda', size=36, color=WHITE, bold=True)

agenda_items = [
    ('01', 'Product Audit', '26 issues identified across 9 critical, 8 high, 9 medium'),
    ('02', 'Competitive Landscape', 'Benchmarked vs. Otter.ai, MacWhisper, Descript, Rev.com'),
    ('03', 'Platform Alignment', 'macOS 26 Tahoe, Liquid Glass, Foundation Models'),
    ('04', 'Feature 1: Per-File Transcription', 'Transcribe individual files on-demand'),
    ('05', 'Feature 2: Bulk Summarization', 'Batch AI summarization with Apple Intelligence'),
    ('06', 'Feature 3: Smart File Renaming', 'AI-powered naming from date, content, speakers'),
    ('07', 'UX Polish & Accessibility', 'Keyboard shortcuts, VoiceOver, design system'),
    ('08', 'Implementation Roadmap', '5-week phased delivery plan'),
]

for i, (num, title, desc) in enumerate(agenda_items):
    y = Inches(1.4) + Inches(i * 0.72)
    card = add_shape(slide, Inches(0.8), y, Inches(11.7), Inches(0.6), BG_CARD, 0.05)
    add_text(slide, Inches(1.1), y + Inches(0.05), Inches(0.6), Inches(0.5),
             num, size=18, color=ACCENT, bold=True)
    add_text(slide, Inches(1.8), y + Inches(0.05), Inches(3), Inches(0.5),
             title, size=18, color=WHITE, bold=True)
    add_text(slide, Inches(5), y + Inches(0.08), Inches(7), Inches(0.5),
             desc, size=14, color=LIGHT_GRAY)

# =====================================================================
# SLIDE 3: PRODUCT AUDIT — CRITICAL ISSUES
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         'Product Audit: Critical Issues', size=36, color=WHITE, bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.5),
         '9 critical/high issues require immediate attention before new feature work',
         size=16, color=LIGHT_GRAY)

issues = [
    ('C1', 'Description edits lost on restart', 'Not persisted to database', RED),
    ('C2', 'Auto-start setting does nothing', 'Stored but never read', RED),
    ('C3', 'Default locale setting ignored', 'Always defaults to en-US', RED),
    ('C5', 'Language field stores URL', 'PersistedJobData bug', RED),
    ('C6', 'Security scope 3s timeout', 'Large files may fail', RED),
    ('C7', 'Topic extraction never called', 'AI code exists, not wired', ORANGE),
    ('C8', 'Speaker detection never called', 'AI code exists, not wired', ORANGE),
    ('C9', 'AI filename never called', 'Fully implemented, not wired', ORANGE),
]

for i, (code, title, detail, color) in enumerate(issues):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + Inches(col * 6.2)
    y = Inches(1.6) + Inches(row * 1.2)
    card = add_shape(slide, x, y, Inches(5.8), Inches(1.0), BG_CARD, 0.05)
    # Color indicator
    indicator = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.15), y + Inches(0.15), Inches(0.08), Inches(0.7))
    indicator.fill.solid()
    indicator.fill.fore_color.rgb = color
    indicator.line.fill.background()
    add_text(slide, x + Inches(0.4), y + Inches(0.1), Inches(0.6), Inches(0.4),
             code, size=14, color=color, bold=True)
    add_text(slide, x + Inches(1.0), y + Inches(0.1), Inches(4.5), Inches(0.4),
             title, size=16, color=WHITE, bold=True)
    add_text(slide, x + Inches(1.0), y + Inches(0.5), Inches(4.5), Inches(0.4),
             detail, size=13, color=LIGHT_GRAY)

# =====================================================================
# SLIDE 4: PRODUCT AUDIT — UX ISSUES
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         'Product Audit: UX & Polish Gaps', size=36, color=WHITE, bold=True)

# Left column: UX Issues
add_text(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(0.5),
         'High-Priority UX Issues', size=20, color=ACCENT, bold=True)

ux_issues = [
    'Blank right pane for pending/processing files',
    'No confirmation for destructive actions',
    'Inconsistent file addition flows',
    'Chat stop button doesn\'t cancel',
    'No onboarding experience',
    'Start button disabled without explanation',
]
add_bullet_list(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4), ux_issues, size=14)

# Right column: Polish Gaps
add_text(slide, Inches(7), Inches(1.2), Inches(5.5), Inches(0.5),
         'Polish & Interaction Gaps', size=20, color=ACCENT, bold=True)

polish_issues = [
    'Zero keyboard shortcuts defined',
    'No context menu on transcript text',
    'Three inconsistent rename entry points',
    'No animations between states',
    'Inconsistent corner radii & hard-coded colors',
    'Debug print statements throughout code',
]
add_bullet_list(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(4), polish_issues, size=14)

# Bottom stat bar
add_shape(slide, Inches(0.8), Inches(6.2), Inches(3.6), Inches(0.8), BG_CARD, 0.05)
add_text(slide, Inches(1.0), Inches(6.25), Inches(3.2), Inches(0.35),
         '9 Critical', size=24, color=RED, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1.0), Inches(6.6), Inches(3.2), Inches(0.3),
         'Data loss / broken features', size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_shape(slide, Inches(4.8), Inches(6.2), Inches(3.6), Inches(0.8), BG_CARD, 0.05)
add_text(slide, Inches(5.0), Inches(6.25), Inches(3.2), Inches(0.35),
         '8 High', size=24, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(5.0), Inches(6.6), Inches(3.2), Inches(0.3),
         'Poor user experience', size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_shape(slide, Inches(8.8), Inches(6.2), Inches(3.6), Inches(0.8), BG_CARD, 0.05)
add_text(slide, Inches(9.0), Inches(6.25), Inches(3.2), Inches(0.35),
         '9 Medium', size=24, color=YELLOW, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(9.0), Inches(6.6), Inches(3.2), Inches(0.3),
         'Polish / completeness', size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 5: COMPETITIVE LANDSCAPE
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         'Competitive Landscape', size=36, color=WHITE, bold=True)

# Comparison cards
competitors = [
    ('Otter.ai', 'Cloud AI', [
        'Auto-summaries after every meeting',
        'Cross-library AI Chat',
        'Speaker naming is weak',
        'English-only transcription',
    ]),
    ('MacWhisper', 'Local Whisper', [
        'Native macOS feel, Liquid Glass adopted',
        'Batch + Watch Folders (auto)',
        'Multi-format export in batch mode',
        'Already $69 lifetime (no subscription)',
    ]),
    ('Descript', 'Cloud', [
        'Text-based media editing paradigm',
        'Strong speaker auto-detection',
        'Pricing penalizes multi-file workflows',
        'Electron wrapper, not native macOS',
    ]),
    ('Rev.com', 'Cloud + Human', [
        'Multi-file AI chat across library',
        'Speaker diarization + label merge',
        '96% AI / 99% human accuracy',
        'Legal-grade encryption',
    ]),
]

for i, (name, model_type, points) in enumerate(competitors):
    x = Inches(0.8) + Inches(i * 3.15)
    card = add_shape(slide, x, Inches(1.3), Inches(2.9), Inches(5.5), BG_CARD, 0.03)
    add_text(slide, x + Inches(0.2), Inches(1.4), Inches(2.5), Inches(0.5),
             name, size=20, color=ACCENT, bold=True)
    add_text(slide, x + Inches(0.2), Inches(1.9), Inches(2.5), Inches(0.3),
             model_type, size=12, color=LIGHT_GRAY)
    add_bullet_list(slide, x + Inches(0.2), Inches(2.4), Inches(2.5), Inches(4),
                    points, size=13, color=WHITE)

# =====================================================================
# SLIDE 6: NITNAB ADVANTAGES
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         'NitNab\'s Competitive Position', size=36, color=WHITE, bold=True)

# Advantages
add_text(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.5),
         'Unique Advantages', size=22, color=GREEN, bold=True)

advantages = [
    'Fully on-device: no cloud, no subscription, complete privacy',
    'Apple Intelligence: native Foundation Models integration',
    'Company context: vocabulary + people databases for accuracy',
    'macOS native: true SwiftUI, not Electron',
    'iCloud sync: structured folder output per transcription',
]
add_bullet_list(slide, Inches(0.8), Inches(1.9), Inches(5.5), Inches(4), advantages, size=15)

# Gaps
add_text(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(0.5),
         'Gaps to Close', size=22, color=RED, bold=True)

gaps = [
    'Per-file transcription (every competitor has this)',
    'Auto-summarization (Otter does this automatically)',
    'Wire speaker identification (code exists, not connected)',
    'Smart naming with AI (no competitor does this well)',
]
add_bullet_list(slide, Inches(7), Inches(1.9), Inches(5.5), Inches(4), gaps, size=15)

# Bottom insight
add_shape(slide, Inches(0.8), Inches(5.8), Inches(11.7), Inches(1.0), BG_CARD, 0.05)
add_text(slide, Inches(1.2), Inches(5.9), Inches(11), Inches(0.8),
         'Key Insight: NitNab has 3 fully-implemented AI features (topics, speakers, filenames) '
         'that are never called. Wiring these up immediately closes the gap with competitors.',
         size=16, color=YELLOW)

# =====================================================================
# SLIDE 7: LIQUID GLASS ALIGNMENT
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         'Platform Alignment: Liquid Glass & macOS 26', size=36, color=WHITE, bold=True)

add_text(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
         'Apple\'s most significant design evolution since iOS 7 (WWDC 2025)',
         size=16, color=LIGHT_GRAY)

# Left: What it is
add_text(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(0.5),
         'What is Liquid Glass?', size=20, color=ACCENT, bold=True)

lg_points = [
    'Translucent material with real-time light bending',
    'Specular highlights responding to device motion',
    'Adaptive shadows and interactive behaviors',
    'Completely transparent menu bar on macOS',
    'Navigation layer only (never apply to content)',
]
add_bullet_list(slide, Inches(0.8), Inches(2.4), Inches(5.5), Inches(3), lg_points, size=14)

# Right: What NitNab must do
add_text(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(0.5),
         'NitNab Action Items', size=20, color=ACCENT, bold=True)

actions = [
    'Recompile with Xcode 26 (auto-adopts for standard controls)',
    'Replace custom VStack sidebar with NavigationSplitView',
    'Remove hard-coded Color.blue.opacity() backgrounds',
    'Use system materials instead of custom backgrounds',
    'Test Light / Dark / Increased Contrast modes',
    'Adopt .glassEffect() for custom floating controls',
]
add_bullet_list(slide, Inches(7), Inches(2.4), Inches(5.5), Inches(3.5), actions, size=14)

# API table
add_shape(slide, Inches(0.8), Inches(5.6), Inches(11.7), Inches(1.4), BG_CARD, 0.03)
add_text(slide, Inches(1.0), Inches(5.7), Inches(2.5), Inches(0.35),
         'Key New APIs:', size=14, color=ACCENT, bold=True)
apis = '.glassEffect()   |   GlassEffectContainer   |   .glassEffectID()   |   .backgroundExtensionEffect()   |   .buttonStyle(.glassProminent)   |   ToolbarSpacer'
add_text(slide, Inches(1.0), Inches(6.1), Inches(11), Inches(0.6),
         apis, size=13, color=LIGHT_GRAY)

# =====================================================================
# SLIDE 8: FEATURE 1 — PER-FILE TRANSCRIPTION
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.8),
         'Feature 1: Per-File Transcription', size=36, color=WHITE, bold=True)

# Tag
add_shape(slide, Inches(10), Inches(0.45), Inches(2.5), Inches(0.45), ACCENT, 0.5)
add_text(slide, Inches(10), Inches(0.45), Inches(2.5), Inches(0.45),
         'TABLE STAKES', size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
         'Every competitor offers this. Currently NitNab can only batch-process all pending files.',
         size=16, color=LIGHT_GRAY)

# Components
components = [
    ('FIFO Queue', 'Sequential processing queue in ViewModel. Per-file adds to front, "Transcribe All" adds all pending.'),
    ('Row Button', 'Inline "Transcribe" button per file row. Animates to progress bar during processing.'),
    ('Pending View', 'Right-pane view for pending files: metadata, locale picker, "Start Transcription" CTA.'),
    ('Auto-Start', 'Wire existing broken setting. Drop a file and transcription begins automatically.'),
    ('Header Update', 'Rename to "Transcribe All (N)" with count badge and current-job subtitle.'),
]

for i, (title, desc) in enumerate(components):
    y = Inches(1.7) + Inches(i * 1.0)
    card = add_shape(slide, Inches(0.8), y, Inches(11.7), Inches(0.85), BG_CARD, 0.03)
    add_text(slide, Inches(1.2), y + Inches(0.08), Inches(2.5), Inches(0.4),
             title, size=18, color=ACCENT, bold=True)
    add_text(slide, Inches(1.2), y + Inches(0.45), Inches(11), Inches(0.4),
             desc, size=14, color=WHITE)

# =====================================================================
# SLIDE 9: FEATURE 2 — BULK SUMMARIZATION
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.8),
         'Feature 2: Bulk Summarization', size=36, color=WHITE, bold=True)

add_shape(slide, Inches(10), Inches(0.45), Inches(2.5), Inches(0.45), GREEN, 0.5)
add_text(slide, Inches(10), Inches(0.45), Inches(2.5), Inches(0.45),
         'DIFFERENTIATOR', size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
         'Otter.ai auto-summarizes. No competitor offers explicit bulk summarize. NitNab can lead here.',
         size=16, color=LIGHT_GRAY)

components2 = [
    ('Data Model', 'Add summary, summaryGeneratedAt, summaryStatus to TranscriptionJob. Database migration.'),
    ('Auto-Summarize', 'Generate summary automatically after each transcription completes. Wire extractTopics() and extractNames().'),
    ('Bulk Action', '"Summarize All (N)" button. Sequential processing with progress indicator.'),
    ('Summary Card', 'Summary displayed ABOVE transcript tabs (not buried in a tab). Collapsible with regenerate option.'),
    ('Tag Cloud', 'Wire extractTopics() to populate the currently-empty Advanced mode tag cloud.'),
]

for i, (title, desc) in enumerate(components2):
    y = Inches(1.7) + Inches(i * 1.0)
    card = add_shape(slide, Inches(0.8), y, Inches(11.7), Inches(0.85), BG_CARD, 0.03)
    add_text(slide, Inches(1.2), y + Inches(0.08), Inches(2.5), Inches(0.4),
             title, size=18, color=ACCENT, bold=True)
    add_text(slide, Inches(1.2), y + Inches(0.45), Inches(11), Inches(0.4),
             desc, size=14, color=WHITE)

# =====================================================================
# SLIDE 10: FEATURE 3 — SMART RENAMING
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.8),
         'Feature 3: Smart File Renaming', size=36, color=WHITE, bold=True)

add_shape(slide, Inches(10), Inches(0.45), Inches(2.5), Inches(0.45), ACCENT, 0.5)
add_text(slide, Inches(10), Inches(0.45), Inches(2.5), Inches(0.45),
         'UNIQUE TO NITNAB', size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
         'No competitor offers AI-powered multi-component naming. User picks their formula, AI fills it in.',
         size=16, color=LIGHT_GRAY)

# Example
add_shape(slide, Inches(0.8), Inches(1.7), Inches(11.7), Inches(1.0), BG_CARD, 0.03)
add_text(slide, Inches(1.2), Inches(1.75), Inches(11), Inches(0.35),
         'Example Output:', size=14, color=LIGHT_GRAY)
add_text(slide, Inches(1.2), Inches(2.1), Inches(11), Inches(0.45),
         '2026-02-15  ---  Q1 Budget Review  ---  Sarah, Mike', size=22, color=GREEN, bold=True)

components3 = [
    ('Naming Preferences', 'Configurable struct: date format, description style, name style, separator, max length.'),
    ('Smart Rename Sheet', 'Live preview with component toggles. User sees the name update in real-time.'),
    ('Post-Transcription', 'Auto-rename ON: apply silently. OFF: show sheet with suggestion. Default: OFF.'),
    ('Bulk Rename', 'Review sheet showing all proposed renames before applying. Never rename silently in bulk.'),
]

for i, (title, desc) in enumerate(components3):
    y = Inches(2.9) + Inches(i * 1.0)
    card = add_shape(slide, Inches(0.8), y, Inches(11.7), Inches(0.85), BG_CARD, 0.03)
    add_text(slide, Inches(1.2), y + Inches(0.08), Inches(2.5), Inches(0.4),
             title, size=18, color=ACCENT, bold=True)
    add_text(slide, Inches(1.2), y + Inches(0.45), Inches(11), Inches(0.4),
             desc, size=14, color=WHITE)

# =====================================================================
# SLIDE 11: UX POLISH
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         'UX Polish & Usability', size=36, color=WHITE, bold=True)

# Keyboard shortcuts
add_text(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.5),
         'Keyboard Shortcuts (0 exist today)', size=20, color=ACCENT, bold=True)

shortcuts = [
    'Cmd+O: Add files',
    'Cmd+Return: Transcribe selected',
    'Cmd+Shift+Return: Transcribe all',
    'Cmd+E: Export',
    'Cmd+F: Search',
    'Cmd+1/2: Toggle mode',
    'Return: Rename file',
    'Cmd+,: Settings',
]
add_bullet_list(slide, Inches(0.8), Inches(1.9), Inches(5), Inches(4), shortcuts, size=13)

# Other polish
add_text(slide, Inches(7), Inches(1.3), Inches(5.5), Inches(0.5),
         'Additional Polish', size=20, color=ACCENT, bold=True)

other_polish = [
    'Confirmation dialogs for all destructive actions',
    'Pending/Processing views for right pane',
    'Unified file addition flow (company picker on all paths)',
    'DesignSystem constants (radii, colors, spacing)',
    'Animations between state transitions',
    'Onboarding / first-run experience',
    'Replace debug prints with os.Logger',
    'Context menus on transcript text',
]
add_bullet_list(slide, Inches(7), Inches(1.9), Inches(5.5), Inches(4.5), other_polish, size=13)

# =====================================================================
# SLIDE 12: ACCESSIBILITY
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         'Accessibility & Inclusivity', size=36, color=WHITE, bold=True)

acc_items = [
    ('VoiceOver Labels', 'Add descriptive accessibilityLabel with status context to all interactive elements'),
    ('Progress Announcements', 'Use AccessibilityNotification.Announcement for transcription status changes'),
    ('Keyboard Navigation', 'Full Tab/arrow key support. File list arrow navigation. Tab bar switching.'),
    ('High Contrast', 'Test all views with Increase Contrast. Replace hard-coded colors with system semantics.'),
    ('Reduce Transparency', 'Verify Liquid Glass readability. Ensure content remains legible.'),
    ('Writing Tools', 'Standard NSTextView/TextKit 2 for transcript gives free Writing Tools support.'),
]

for i, (title, desc) in enumerate(acc_items):
    y = Inches(1.3) + Inches(i * 0.95)
    card = add_shape(slide, Inches(0.8), y, Inches(11.7), Inches(0.8), BG_CARD, 0.03)
    add_text(slide, Inches(1.2), y + Inches(0.05), Inches(3), Inches(0.35),
             title, size=16, color=ACCENT, bold=True)
    add_text(slide, Inches(1.2), y + Inches(0.4), Inches(11), Inches(0.35),
             desc, size=13, color=WHITE)

# =====================================================================
# SLIDE 13: IMPLEMENTATION ROADMAP
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
         'Implementation Roadmap', size=36, color=WHITE, bold=True)

phases = [
    ('Week 1', 'Foundation & Bug Fixes', '~16 hrs', 'Fix 9 critical bugs, add confirmations, unify flows, create design system', RED),
    ('Week 2', 'Per-File Transcription', '~16 hrs', 'FIFO queue, per-row button, pending/processing views, auto-start setting', ACCENT),
    ('Week 3', 'Bulk Summarization', '~16 hrs', 'Summary model, DB migration, auto-summarize, summary card UI, tag cloud', GREEN),
    ('Week 4', 'Smart File Renaming', '~17 hrs', 'Naming model, AI generation, rename sheet, settings, bulk rename', ACCENT),
    ('Week 5', 'Polish & Accessibility', '~21 hrs', 'Liquid Glass, keyboard shortcuts, VoiceOver, animations, onboarding', LIGHT_GRAY),
]

for i, (week, title, hours, desc, color) in enumerate(phases):
    y = Inches(1.2) + Inches(i * 1.15)

    # Week indicator
    indicator = add_shape(slide, Inches(0.8), y, Inches(1.2), Inches(1.0), color, 0.08)
    add_text(slide, Inches(0.8), y + Inches(0.15), Inches(1.2), Inches(0.35),
             week, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.8), y + Inches(0.5), Inches(1.2), Inches(0.3),
             hours, size=12, color=WHITE, align=PP_ALIGN.CENTER)

    # Content card
    card = add_shape(slide, Inches(2.2), y, Inches(10.3), Inches(1.0), BG_CARD, 0.03)
    add_text(slide, Inches(2.5), y + Inches(0.1), Inches(9.5), Inches(0.4),
             title, size=20, color=WHITE, bold=True)
    add_text(slide, Inches(2.5), y + Inches(0.55), Inches(9.5), Inches(0.4),
             desc, size=14, color=LIGHT_GRAY)

# Total
add_shape(slide, Inches(0.8), Inches(7.0), Inches(11.7), Inches(0.01), ACCENT)

# =====================================================================
# SLIDE 14: NEXT STEPS
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
         'Next Steps', size=42, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

next_steps = [
    '1.  Approve plan and prioritize phases',
    '2.  Fix critical bugs (Week 1) before adding features',
    '3.  Begin Per-File Transcription implementation (Week 2)',
    '4.  Wire disconnected AI functions alongside Feature 2',
    '5.  Adopt Liquid Glass design system in Phase 5',
]

for i, step in enumerate(next_steps):
    y = Inches(2.8) + Inches(i * 0.7)
    add_text(slide, Inches(2), y, Inches(9), Inches(0.6),
             step, size=20, color=WHITE if i > 0 else ACCENT)

add_text(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
         'Design & Engineering Review Team  |  Confidential',
         size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# =====================================================================
# SAVE
# =====================================================================
prs.save('/home/user/nitnab/docs/NitNab_Design_Review_and_Feature_Plan.pptx')
print('PPTX generated successfully.')
